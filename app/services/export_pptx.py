"""
Export PPTX pixel-perfect depuis un payload slides_5_0.

Conversion : px (canvas 1920x1080) -> EMU (English Metric Units)
  1 px a 96 dpi -> 1 pt -> 12700 EMU
  PPTX 16:9 standard = 12192000 x 6858000 EMU
  Scale X = 12192000 / 1920 = 6350 EMU/px
  Scale Y = 6858000 / 1080 = 6350 EMU/px
"""

from __future__ import annotations

import io
import logging
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from app.api.schemas.common import Study
from app.schemas.slide_objects import validate_slide_objects
from app.services.slides_5_0_builder import build_slides_5_0_for_study

logger = logging.getLogger(__name__)

_SLIDE_W_EMU = 12192000
_SLIDE_H_EMU = 6858000
_CANVAS_W_PX = 1920
_CANVAS_H_PX = 1080
_PX_TO_EMU = _SLIDE_W_EMU / _CANVAS_W_PX  # 6350.0

_DARK_BG = RGBColor(0x0A, 0x10, 0x1A)
_LIGHT_BG = RGBColor(0xFF, 0xFF, 0xFF)
_DEFAULT_TEXT_DARK = RGBColor(0xE2, 0xE8, 0xF0)
_DEFAULT_TEXT_LIGHT = RGBColor(0x1E, 0x29, 0x3B)


def _px(value, fallback=0):
    if value is None:
        return Emu(fallback * _PX_TO_EMU)
    return Emu(int(value * _PX_TO_EMU))


def _hex_to_rgb(hex_color):
    if not hex_color:
        return None
    h = hex_color.lstrip("#")
    if len(h) != 6:
        return None
    try:
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except ValueError:
        return None


def _apply_run_style(run, style, is_dark_bg):
    font = run.font
    font_size_px = style.get("fontSize") or style.get("font_size")
    if font_size_px:
        try:
            font.size = Pt(float(str(font_size_px).replace("px", "").strip()) * 0.75)
        except (ValueError, TypeError):
            pass
    color_val = style.get("color") or style.get("colour")
    rgb = _hex_to_rgb(str(color_val)) if color_val else None
    if rgb:
        font.color.rgb = rgb
    else:
        font.color.rgb = _DEFAULT_TEXT_DARK if is_dark_bg else _DEFAULT_TEXT_LIGHT
    font_weight = style.get("fontWeight") or style.get("font_weight", "")
    if str(font_weight) in ("700", "bold"):
        font.bold = True
    if style.get("fontStyle") == "italic":
        font.italic = True


def _add_text_box(slide, obj, is_dark_bg):
    text = obj.get("text") or ""
    style = obj.get("style") or {}
    txBox = slide.shapes.add_textbox(
        _px(obj.get("left")), _px(obj.get("top")),
        _px(obj.get("width", 200)), _px(obj.get("height", 40))
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    children = obj.get("children") or []
    if children:
        for i, child in enumerate(children):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = child.get("text") or ""
            _apply_run_style(run, child.get("style") or {}, is_dark_bg)
    else:
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        _apply_run_style(run, style, is_dark_bg)


def _add_shape(slide, obj, css_vars, is_dark_bg):
    left = _px(obj.get("left") or 0)
    top = _px(obj.get("top") or 0)
    width = _px(obj.get("width") or 100)
    height = _px(obj.get("height") or 100)
    style = obj.get("style") or {}
    bg_color_raw = style.get("backgroundColor") or style.get("background")
    if bg_color_raw and str(bg_color_raw).startswith("var("):
        var_name = str(bg_color_raw)[4:-1]
        bg_color_raw = css_vars.get(var_name, "#334155")
    rgb = _hex_to_rgb(str(bg_color_raw)) if bg_color_raw else None
    shape = slide.shapes.add_shape(1, left, top, width, height)
    try:
        shape.line.fill.background()
    except Exception:
        try:
            shape.line.width = Pt(0)
        except Exception:
            pass
    if rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb
    else:
        try:
            shape.fill.background()
        except Exception:
            pass
    children = obj.get("children") or []
    if children:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = _px(16)
        tf.margin_top = _px(12)
        tf.margin_right = _px(16)
        tf.margin_bottom = _px(12)
        first_para = True
        for child in children:
            child_text = child.get("text") or ""
            if not child_text:
                continue
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
            run = p.add_run()
            run.text = child_text
            _apply_run_style(run, child.get("style") or {}, is_dark_bg)
            try:
                p.space_after = Pt(4)
            except Exception:
                pass


def _add_native_pie_chart(slide, categories, values, left_emu, top_emu, width_emu, height_emu, title=""):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Repartition", values)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left_emu, top_emu, width_emu, height_emu, chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = True
    try:
        chart.legend.position = XL_LEGEND_POSITION.RIGHT
        chart.legend.include_in_layout = False
    except Exception:
        pass
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title


def _add_native_bar_chart(slide, categories, values, left_emu, top_emu, width_emu, height_emu, title=""):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Score", values)
    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left_emu, top_emu, width_emu, height_emu, chart_data
    )
    chart = graphic_frame.chart
    chart.has_legend = False
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title


def _add_chart_native(slide, obj):
    chart_data = obj.get("chart_data") or {}
    chart_type = chart_data.get("type", "bar")
    categories = chart_data.get("categories") or []
    values = chart_data.get("values") or []
    title = chart_data.get("title", "")
    if not categories or not values or len(categories) != len(values):
        logger.warning("[pptx] chart_native %s : donnees manquantes", obj.get("id"))
        return
    try:
        values = [float(v) for v in values]
    except (TypeError, ValueError):
        logger.warning("[pptx] chart_native %s : valeurs non numeriques", obj.get("id"))
        return
    left = int(_px(obj.get("left")))
    top = int(_px(obj.get("top")))
    width = int(_px(obj.get("width", 400)))
    height = int(_px(obj.get("height", 300)))
    if chart_type == "pie":
        _add_native_pie_chart(slide, categories, values, left, top, width, height, title)
    else:
        _add_native_bar_chart(slide, categories, values, left, top, width, height, title)


def _apply_slide_background(prs_slide, background, css_vars):
    from pptx.oxml.ns import qn
    import lxml.etree as etree
    is_dark = background == "dark"
    bg_hex = css_vars.get("--stella-bg", "#0A101A") if is_dark else "#FFFFFF"
    rgb = _hex_to_rgb(bg_hex) or (_DARK_BG if is_dark else _LIGHT_BG)
    bg = prs_slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def build_pptx_for_study(study: Study) -> bytes:
    payload = build_slides_5_0_for_study(study)
    slides_data = payload.get("slides", [])
    css_vars_dict = (payload.get("css_vars") or {}).get("variables") or {}

    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W_EMU)
    prs.slide_height = Emu(_SLIDE_H_EMU)
    blank_layout = prs.slide_layouts[6]

    for slide_info in slides_data:
        prs_slide = prs.slides.add_slide(blank_layout)
        background = slide_info.get("background", "dark")
        is_dark = background == "dark"
        _apply_slide_background(prs_slide, background, css_vars_dict)

        raw_objects = slide_info.get("objects") or []
        objects = validate_slide_objects(raw_objects)

        for obj in objects:
            obj_type = obj.get("data_object_type", "textbox")
            try:
                if obj_type == "textbox":
                    if not obj.get("text") and not obj.get("children"):
                        continue
                    _add_text_box(prs_slide, obj, is_dark)
                elif obj_type == "shape":
                    _add_shape(prs_slide, obj, css_vars_dict, is_dark)
                elif obj_type == "kpi_list":
                    items = obj.get("items") or []
                    lines = [
                        f"{it.get('label', '')}: {it.get('value', '')}"
                        for it in items if it.get("label") or it.get("value")
                    ]
                    _add_text_box(prs_slide, {**obj, "text": "\n".join(lines), "children": []}, is_dark)
                elif obj_type == "chart_native":
                    _add_chart_native(prs_slide, obj)
                elif obj_type in ("chart", "image", "icon"):
                    _add_text_box(prs_slide, {**obj, "text": f"[{obj_type.upper()}]", "children": []}, is_dark)
            except Exception as exc:
                logger.warning("[pptx] objet %s ignore : %s", obj.get("id"), exc)

        notes_slide = prs_slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = f"{slide_info.get('section_id', '')} -- {slide_info.get('title', '')}"

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
